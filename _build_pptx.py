"""Genera la presentacion del TP1."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
DARK = RGBColor(44, 62, 80)
ACCENT = RGBColor(108, 52, 131)  # #6C3483
LIGHT_BG = RGBColor(245, 245, 250)
BLUE = RGBColor(41, 128, 185)
RED = RGBColor(192, 57, 43)

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=DARK, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=DARK):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(6)
        p.level = 0
    return tf

def add_image(slide, path, left, top, width=None, height=None):
    if os.path.exists(path):
        if width and height:
            slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(width), Inches(height))
        elif width:
            slide.shapes.add_picture(path, Inches(left), Inches(top), width=Inches(width))
        elif height:
            slide.shapes.add_picture(path, Inches(left), Inches(top), height=Inches(height))

def title_bar(slide, title, subtitle=None):
    """Barra de titulo en la parte superior."""
    shape = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(13.333), Inches(1.2)  # 1 = rectangle
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()

    add_textbox(slide, 0.5, 0.15, 12, 0.7, title, font_size=32, bold=True, color=WHITE)
    if subtitle:
        add_textbox(slide, 0.5, 0.75, 12, 0.4, subtitle, font_size=16, color=RGBColor(220, 220, 230))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Portada
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, ACCENT)

add_textbox(slide, 1, 1.5, 11, 1.5,
            "TP1: Regresion e Introduccion\na la evaluacion de modelos",
            font_size=40, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 1, 3.5, 11, 0.8,
            "72.75 Aprendizaje Automatico -- ITBA -- 2026 Q1",
            font_size=22, color=RGBColor(220, 210, 230), alignment=PP_ALIGN.CENTER)

add_textbox(slide, 1, 4.5, 11, 0.8,
            "Dataset: Wine Quality (Cortez et al., 2009)",
            font_size=20, color=RGBColor(200, 190, 210), alignment=PP_ALIGN.CENTER)

add_textbox(slide, 1, 6.0, 11, 0.5,
            "Fecha de defensa: 25/03/2026",
            font_size=16, color=RGBColor(180, 170, 195), alignment=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Intro teorica: train/val/test
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
title_bar(slide, "Introduccion teorica", "Separacion train - validacion - test")

add_bullet_list(slide, 0.5, 1.5, 5.8, 5.5, [
    "Por que separar los datos?",
    "",
    "TRAIN (64%): el modelo aprende los patrones.",
    "",
    "VALIDATION (16%): se usa para ajustar hiperparametros",
    "  (grado polinomio, alpha de Lasso) y seleccionar el",
    "  mejor modelo. Se implementa via K-Fold CV (k=5).",
    "",
    "TEST (20%): se evalua UNA sola vez al final.",
    "  Estima el error real en datos nunca vistos.",
    "",
    "Si usaramos test para elegir el modelo, estariamos",
    "optimizando sobre datos que deberian ser 'no vistos',",
    "invalidando la estimacion de error (data leakage).",
], font_size=15)

# Diagrama textual
add_textbox(slide, 6.8, 1.5, 6, 5.5,
    "Dataset limpio (6456 filas)\n"
    "          |\n"
    "    80% Train (5164)\n"
    "          |\n"
    "     K-Fold CV (k=5)\n"
    "     estratificado\n"
    "          |\n"
    "    Fold 1: train / val\n"
    "    Fold 2: train / val\n"
    "    Fold 3: train / val\n"
    "    Fold 4: train / val\n"
    "    Fold 5: train / val\n"
    "\n"
    "    20% Test (1292)\n"
    "    Solo evaluacion final\n"
    "\n"
    "Stratified: mantiene la proporcion\n"
    "de cada quality (3-9) en cada split.",
    font_size=14, color=DARK, font_name="Consolas")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Dataset
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
title_bar(slide, "El dataset", "Wine Quality -- Vinho Verde, Portugal")

add_bullet_list(slide, 0.5, 1.5, 5.5, 3, [
    "6497 muestras (1599 tintos + 4898 blancos)",
    "11 variables fisicoquimicas de entrada",
    "1 variable target: quality (0-10, evaluada por expertos)",
    "Sin valores faltantes",
    "",
    "Variables: fixed acidity, volatile acidity, citric acid,",
    "residual sugar, chlorides, free/total sulfur dioxide,",
    "density, pH, sulphates, alcohol",
], font_size=15)

add_image(slide, "results/tipo_vino.png", 6.5, 1.3, height=2.8)
add_image(slide, "results/distribucion_quality.png", 6.5, 4.2, height=2.8)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Variable categorica
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
title_bar(slide, "1.1 Variable categorica + 1.2 Valores faltantes")

add_bullet_list(slide, 0.5, 1.5, 6, 2.5, [
    "Variable categorica: tipo de vino (tinto / blanco)",
    "",
    "Estrategia: One-Hot Encoding con drop_first=True",
    "  Se genera una columna: red_wine_type (1=tinto, 0=blanco)",
    "  No se necesita white_wine_type (seria 1 - red_wine_type)",
    "  Mantener ambas introduce multicolinealidad perfecta",
    "",
    "Valores faltantes: no hay ninguno (verificado con isnull().sum())",
], font_size=16)

add_bullet_list(slide, 0.5, 4.5, 12, 2.5, [
    "Por que incluir el tipo de vino como feature?",
    "  Las distribuciones difieren significativamente entre tipos:",
    "  volatile acidity: media 0.53 (tinto) vs 0.28 (blanco)",
    "  total sulfur dioxide: media 46 (tinto) vs 138 (blanco)",
    "  Incluirlo permite al modelo capturar estas diferencias sistematicas.",
], font_size=15)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Outliers: distribucion
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
title_bar(slide, "1.3 Outliers -- Analisis de distribucion")

add_image(slide, "data/eda/distribucion_combined.png", 0.2, 1.3, width=6.4)
add_image(slide, "data/eda/boxplot_combined.png", 6.7, 1.3, width=6.4)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Outliers: estrategia
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
title_bar(slide, "1.3 Outliers -- Estrategia de limpieza")

add_textbox(slide, 0.5, 1.4, 6, 0.5, "Eliminar (41 filas, 0.63% del dataset):",
            font_size=18, bold=True, color=ACCENT)

add_bullet_list(slide, 0.5, 1.9, 6, 3, [
    "total sulfur dioxide > 300 mg/L (6 filas) -- supera limite legal UE",
    "free sulfur dioxide > 150 mg/L (1 fila) -- fisicamente incoherente",
    "chlorides > 0.3 g/L (24 filas) -- 3x-6x el maximo normal",
    "citric acid > 1.0 g/L (2 filas) -- supera limite legal de adicion",
    "density > 1.010 (3 filas) -- incompatible con Vinho Verde",
    "pH < 2.80 (7 filas) -- acidez comparable al vinagre",
], font_size=14)

add_textbox(slide, 0.5, 4.5, 6, 0.5, "Winsorizar al p1-p99 (ajustado solo sobre train):",
            font_size=18, bold=True, color=ACCENT)

add_bullet_list(slide, 0.5, 5.0, 6, 2, [
    "residual sugar, volatile acidity, sulphates, chlorides",
    "Preserva el tamano muestral y la forma de la distribucion",
    "Solo aplasta las colas extremas sin eliminar filas",
], font_size=14)

add_textbox(slide, 7, 1.4, 5.8, 0.5, "Data leakage:",
            font_size=18, bold=True, color=RED)

add_bullet_list(slide, 7, 1.9, 5.8, 5, [
    "Eliminacion (umbrales fijos del dominio):",
    "  No depende de estadisticas -> se aplica ANTES del split",
    "",
    "Winsorizacion (percentiles calculados):",
    "  Depende de estadisticas -> se aplica DESPUES del split",
    "  fit en train, transform en train y test",
    "",
    "Escalado (media/std):",
    "  Depende de estadisticas -> se aplica DESPUES del split",
    "  fit en train, transform en train y test",
    "",
    "Stratified split:",
    "  Solo asigna filas -> NO es data leakage",
], font_size=14)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Caracteristicas: correlacion + escalado
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
title_bar(slide, "1.4 Caracteristicas -- Correlacion y escalado")

add_image(slide, "results/heatmap_correlacion.png", 0.2, 1.3, height=5.8)

add_bullet_list(slide, 7, 1.5, 5.8, 2.5, [
    "Ningun par supera |r| > 0.80",
    "  -> No descartamos variables por correlacion",
    "",
    "Correlaciones mas fuertes con quality:",
    "  alcohol: +0.44  (la mas importante)",
    "  density: -0.31",
    "  volatile acidity: -0.27",
], font_size=15)

add_textbox(slide, 7, 4.2, 5.8, 0.4, "Escalado: StandardScaler",
            font_size=18, bold=True, color=ACCENT)

add_bullet_list(slide, 7, 4.7, 5.8, 2.5, [
    "Ratio de rangos: 16,627x (density vs total SO2)",
    "Obligatorio para Lasso (L1 penaliza por magnitud)",
    "Obligatorio para regresion polinomica (evitar overflow)",
    "fit solo sobre train, transform sobre train y test",
    "Se aplica dentro del Pipeline de sklearn para",
    "  evitar leakage dentro del K-Fold CV",
], font_size=14)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Regresion lineal
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
title_bar(slide, "2. Regresion lineal", "K-Fold CV (k=5) estratificado")

add_image(slide, "results/regresion_lineal.png", 0.3, 1.3, width=12.7)

add_bullet_list(slide, 0.5, 5.5, 12, 1.5, [
    "RMSE val = 0.7329 +/- 0.0154   |   R2 val = 0.294   |   Gap train-val = 0.001 (sin overfitting)",
    "Las variables mas importantes (por coeficiente estandarizado): density (-0.39), residual sugar (+0.32), volatile acidity (-0.23), alcohol (+0.22)",
], font_size=15)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Forward selection
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
title_bar(slide, "2. Regresion lineal -- Forward Selection (Wrapper)")

add_image(slide, "results/forward_selection.png", 0.3, 1.3, width=12.7)

add_bullet_list(slide, 0.5, 5.5, 12, 1.5, [
    "Punto optimo: 10 features (descarta citric acid y chlorides) -- RMSE val = 0.7328",
    "Mejora despreciable vs todas las features (0.0001) -- con 12 variables no hay ruido significativo",
    "Orden de importancia: alcohol > volatile acidity > sulphates > residual sugar > free SO2 > total SO2 > density > wine type > fixed acidity > pH",
], font_size=14)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Lasso
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
title_bar(slide, "3. Regresion polinomica + Lasso (Embedded)")

add_image(slide, "results/regresion_lasso.png", 0.3, 1.3, width=12.7)

add_bullet_list(slide, 0.5, 5.3, 12, 2, [
    "Mejor modelo: Polinomio grado 2 + Lasso (alpha=0.001) -- RMSE val = 0.7003 (+4.4% vs lineal)",
    "Lasso elimina 14 de 90 coeficientes automaticamente (feature selection embedded)",
    "Grado 3 sin regularizacion: overfitting severo (RMSE train=0.63, val=0.76). Lasso lo rescata a 0.70 eliminando 335/454 coefs",
], font_size=14)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Evaluacion final
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
title_bar(slide, "4-5. Evaluacion final y comparacion")

add_image(slide, "results/evaluacion_final.png", 0.15, 1.25, width=13)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Conclusiones / respuestas
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
title_bar(slide, "Conclusiones")

add_textbox(slide, 0.5, 1.5, 12, 0.5,
            "1. Que modelo obtuvo menor error?",
            font_size=20, bold=True, color=ACCENT)
add_bullet_list(slide, 0.8, 2.0, 11, 1, [
    "Polinomio grado 2 + Lasso (alpha=0.001), seleccionado por menor RMSE en cross-validation: 0.7003",
    "Confirmado sobre test set: RMSE = 0.6926 (gap de solo 0.008)",
], font_size=16)

add_textbox(slide, 0.5, 3.1, 12, 0.5,
            "2. Cual elegirian para una aplicacion real?",
            font_size=20, bold=True, color=ACCENT)
add_bullet_list(slide, 0.8, 3.6, 11, 1.5, [
    "Polinomio grado 2 + Lasso. Mejor RMSE en CV, gap bajo con test, Lasso reduce",
    "complejidad (77/90 coefs activos) y la regularizacion lo hace robusto ante datos nuevos.",
], font_size=16)

add_textbox(slide, 0.5, 4.7, 12, 0.5,
            "3. Que RMSE esperan en datos nuevos?",
            font_size=20, bold=True, color=ACCENT)
add_bullet_list(slide, 0.8, 5.2, 11, 2, [
    "RMSE esperado: 0.7003 +/- 0.0117 (estimacion por cross-validation, no por test)",
    "El modelo se equivoca en promedio ~0.7 puntos en la escala de calidad (0-10)",
    "R2 ~ 0.35: el modelo explica 35% de la varianza. La calidad del vino depende de",
    "factores subjetivos (preferencia del catador) que las variables fisicoquimicas no capturan.",
], font_size=16)

# ══════════════════════════════════════════════════════════════════════════════

prs.save("tp1_presentacion.pptx")
print("Presentacion guardada: tp1_presentacion.pptx")
print(f"Total de slides: {len(prs.slides)}")
